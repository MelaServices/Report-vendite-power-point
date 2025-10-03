#!/usr/bin/env python3
"""
Amazon Orders Analyzer - Stile Template Mela
Crea una presentazione PowerPoint usando il template Mela come base.
"""

import os
import glob
from datetime import datetime
from collections import defaultdict
from typing import List, Dict, Any
import csv

# Importa librerie per PowerPoint
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️  python-pptx non disponibile")


class AmazonOrdersAnalyzer:
    """Classe per analizzare gli ordini Amazon."""

    # Colori stile Mela (eleganti e minimali)
    COLORS = {
        'primary': RGBColor(255, 111, 67),     # Arancione/rosso melone
        'secondary': RGBColor(89, 131, 252),    # Blu chiaro
        'accent1': RGBColor(255, 193, 112),     # Giallo/arancio pastello
        'accent2': RGBColor(132, 217, 205),     # Verde acqua
        'accent3': RGBColor(255, 159, 124),     # Pesca
        'text_dark': RGBColor(51, 51, 51),      # Grigio scuro
        'text_light': RGBColor(102, 102, 102),  # Grigio medio
        'bg_light': RGBColor(250, 250, 250),    # Bianco sporco
        'white': RGBColor(255, 255, 255),
    }

    def __init__(self, data_directory: str = "."):
        self.data_directory = data_directory
        self.orders = []
        self.headers = []

    def load_files(self, pattern: str = "*.txt"):
        """Carica tutti i file che corrispondono al pattern."""
        file_pattern = os.path.join(self.data_directory, pattern)
        files = glob.glob(file_pattern)
        files = [f for f in files if not f.endswith("file.txt")]

        print(f"Trovati {len(files)} file da analizzare")

        for file_path in files:
            self._load_file(file_path)

        print(f"Caricati {len(self.orders)} ordini in totale")

    def _load_file(self, file_path: str):
        """Carica un singolo file di ordini."""
        print(f"Caricamento: {os.path.basename(file_path)}")

        with open(file_path, 'r', encoding='utf-8') as f:
            reader = csv.DictReader(f, delimiter='\t')

            if not self.headers:
                self.headers = reader.fieldnames

            for row in reader:
                if row.get('amazon-order-id'):
                    self.orders.append(row)

    def get_summary(self) -> Dict[str, Any]:
        """Ottiene un riepilogo generale degli ordini."""
        summary = {
            'total_orders': len(self.orders),
            'by_status': defaultdict(int),
            'by_marketplace': defaultdict(int),
            'by_country': defaultdict(int),
            'by_fulfillment': defaultdict(int),
            'total_revenue': 0.0,
            'total_tax': 0.0,
            'total_shipping': 0.0,
            'total_items_sold': 0,
        }

        for order in self.orders:
            status = order.get('order-status', 'Unknown')
            summary['by_status'][status] += 1

            marketplace = order.get('sales-channel', 'Unknown')
            summary['by_marketplace'][marketplace] += 1

            country = order.get('ship-country', 'Unknown')
            summary['by_country'][country] += 1

            fulfillment = order.get('fulfillment-channel', 'Unknown')
            summary['by_fulfillment'][fulfillment] += 1

            if status == 'Shipped':
                try:
                    item_price = float(order.get('item-price', 0) or 0)
                    summary['total_revenue'] += item_price

                    item_tax = float(order.get('item-tax', 0) or 0)
                    summary['total_tax'] += item_tax

                    shipping_price = float(order.get('shipping-price', 0) or 0)
                    summary['total_shipping'] += shipping_price
                except (ValueError, TypeError):
                    pass

            try:
                quantity = int(order.get('quantity', 0) or 0)
                summary['total_items_sold'] += quantity
            except (ValueError, TypeError):
                pass

        return summary

    def get_products_summary(self) -> List[Dict[str, Any]]:
        """Ottiene un riepilogo per prodotto."""
        products = defaultdict(lambda: {
            'name': '',
            'sku': '',
            'total_quantity': 0,
            'total_revenue': 0.0,
            'orders_count': 0,
        })

        for order in self.orders:
            sku = order.get('sku', 'Unknown')
            status = order.get('order-status', '')

            if not sku or sku == 'Unknown':
                continue

            if not products[sku]['name']:
                products[sku]['name'] = order.get('product-name', 'N/A')

            products[sku]['sku'] = sku
            products[sku]['orders_count'] += 1

            try:
                quantity = int(order.get('quantity', 0) or 0)
                products[sku]['total_quantity'] += quantity
            except (ValueError, TypeError):
                pass

            if status == 'Shipped':
                try:
                    item_price = float(order.get('item-price', 0) or 0)
                    products[sku]['total_revenue'] += item_price
                except (ValueError, TypeError):
                    pass

        products_list = list(products.values())
        products_list.sort(key=lambda x: x['total_revenue'], reverse=True)

        return products_list

    def get_monthly_trends(self) -> Dict[str, Any]:
        """Analizza i trend mensili degli ordini."""
        monthly = defaultdict(lambda: {
            'orders': 0,
            'revenue': 0.0,
            'items': 0,
            'shipped': 0,
            'cancelled': 0,
            'pending': 0,
        })

        for order in self.orders:
            purchase_date = order.get('purchase-date', '')
            status = order.get('order-status', '')

            if not purchase_date:
                continue

            try:
                dt = datetime.fromisoformat(purchase_date.replace('Z', '+00:00'))
                month_key = dt.strftime('%Y-%m')

                monthly[month_key]['orders'] += 1

                if status == 'Shipped':
                    monthly[month_key]['shipped'] += 1
                elif status == 'Cancelled':
                    monthly[month_key]['cancelled'] += 1
                elif status == 'Pending':
                    monthly[month_key]['pending'] += 1

                quantity = int(order.get('quantity', 0) or 0)
                monthly[month_key]['items'] += quantity

                if status == 'Shipped':
                    item_price = float(order.get('item-price', 0) or 0)
                    monthly[month_key]['revenue'] += item_price

            except (ValueError, TypeError):
                continue

        return dict(sorted(monthly.items()))

    def get_fba_fbm_analysis(self) -> Dict[str, Any]:
        """Analizza differenze tra FBA (Amazon) e FBM (Merchant)."""
        analysis = {
            'fba': {'orders': 0, 'revenue': 0.0, 'items': 0, 'cancelled': 0},
            'fbm': {'orders': 0, 'revenue': 0.0, 'items': 0, 'cancelled': 0},
            'monthly_fba': defaultdict(lambda: {'orders': 0, 'revenue': 0.0}),
            'monthly_fbm': defaultdict(lambda: {'orders': 0, 'revenue': 0.0}),
        }

        for order in self.orders:
            fulfillment = order.get('fulfillment-channel', 'Unknown')
            status = order.get('order-status', '')
            purchase_date = order.get('purchase-date', '')

            # Determina se è FBA o FBM
            is_fba = fulfillment == 'Amazon'
            category = 'fba' if is_fba else 'fbm'

            # Conta ordini
            analysis[category]['orders'] += 1

            # Conta cancellati
            if status == 'Cancelled':
                analysis[category]['cancelled'] += 1

            # Items
            try:
                quantity = int(order.get('quantity', 0) or 0)
                analysis[category]['items'] += quantity
            except (ValueError, TypeError):
                pass

            # Revenue (solo spediti)
            if status == 'Shipped':
                try:
                    item_price = float(order.get('item-price', 0) or 0)
                    analysis[category]['revenue'] += item_price

                    # Analisi mensile
                    if purchase_date:
                        dt = datetime.fromisoformat(purchase_date.replace('Z', '+00:00'))
                        month_key = dt.strftime('%Y-%m')

                        if is_fba:
                            analysis['monthly_fba'][month_key]['orders'] += 1
                            analysis['monthly_fba'][month_key]['revenue'] += item_price
                        else:
                            analysis['monthly_fbm'][month_key]['orders'] += 1
                            analysis['monthly_fbm'][month_key]['revenue'] += item_price
                except (ValueError, TypeError):
                    pass

        # Ordina dati mensili
        analysis['monthly_fba'] = dict(sorted(analysis['monthly_fba'].items()))
        analysis['monthly_fbm'] = dict(sorted(analysis['monthly_fbm'].items()))

        return analysis

    def get_promotions_analysis(self) -> Dict[str, Any]:
        """Analizza le promozioni utilizzate negli ordini (da colonna promotion-ids)."""
        analysis = {
            'total_with_promotions': 0,
            'by_promotion_type': defaultdict(lambda: {'count': 0, 'skus': set()}),
            'skus_on_promotion': defaultdict(int),
            'monthly_promotions': defaultdict(int),
        }

        for order in self.orders:
            promo_id = order.get('promotion-ids', '').strip()
            sku = order.get('sku', 'Unknown')
            purchase_date = order.get('purchase-date', '')

            # Solo ordini con promotion-ids
            if promo_id:
                analysis['total_with_promotions'] += 1

                # Analisi per tipo di promozione
                analysis['by_promotion_type'][promo_id]['count'] += 1
                analysis['by_promotion_type'][promo_id]['skus'].add(sku)

                # SKU in promozione
                analysis['skus_on_promotion'][sku] += 1

                # Analisi mensile
                if purchase_date:
                    try:
                        dt = datetime.fromisoformat(purchase_date.replace('Z', '+00:00'))
                        month_key = dt.strftime('%Y-%m')
                        analysis['monthly_promotions'][month_key] += 1
                    except (ValueError, TypeError):
                        pass

        # Converte sets in count
        for promo_id in analysis['by_promotion_type']:
            analysis['by_promotion_type'][promo_id]['unique_skus'] = len(analysis['by_promotion_type'][promo_id]['skus'])
            # Converte set in lista per facilità d'uso
            analysis['by_promotion_type'][promo_id]['skus'] = list(analysis['by_promotion_type'][promo_id]['skus'])

        # Ordina e converte
        analysis['by_promotion_type'] = dict(analysis['by_promotion_type'])
        analysis['skus_on_promotion'] = dict(analysis['skus_on_promotion'])
        analysis['monthly_promotions'] = dict(sorted(analysis['monthly_promotions'].items()))

        return analysis

    def get_market_monthly_analysis(self) -> Dict[str, Any]:
        """Analizza le vendite mensili per mercato (paese)."""
        analysis = defaultdict(lambda: defaultdict(float))

        for order in self.orders:
            if order.get('order-status') != 'Shipped':
                continue

            country = order.get('ship-country', 'Unknown')
            purchase_date = order.get('purchase-date', '')

            if purchase_date:
                try:
                    dt = datetime.fromisoformat(purchase_date.replace('Z', '+00:00'))
                    month_key = dt.strftime('%Y-%m')

                    item_price = float(order.get('item-price', 0) or 0)
                    analysis[country][month_key] += item_price
                except (ValueError, TypeError):
                    pass

        # Converti in dizionario normale e ordina i mesi
        result = {}
        for country, months in analysis.items():
            result[country] = dict(sorted(months.items()))

        return result

    def get_business_analysis(self) -> Dict[str, Any]:
        """Analizza differenze tra ordini Business e Consumer."""
        analysis = {
            'business': {'orders': 0, 'revenue': 0.0, 'items': 0, 'cancelled': 0},
            'consumer': {'orders': 0, 'revenue': 0.0, 'items': 0, 'cancelled': 0},
            'monthly_business': defaultdict(lambda: {'orders': 0, 'revenue': 0.0}),
            'monthly_consumer': defaultdict(lambda: {'orders': 0, 'revenue': 0.0}),
        }

        for order in self.orders:
            is_business = order.get('is-business-order', '').strip().lower() == 'true'
            status = order.get('order-status', '')
            purchase_date = order.get('purchase-date', '')

            category = 'business' if is_business else 'consumer'

            # Conta ordini
            analysis[category]['orders'] += 1

            # Conta cancellati
            if status == 'Cancelled':
                analysis[category]['cancelled'] += 1

            # Items
            try:
                quantity = int(order.get('quantity', 0) or 0)
                analysis[category]['items'] += quantity
            except (ValueError, TypeError):
                pass

            # Revenue (solo spediti)
            if status == 'Shipped':
                try:
                    item_price = float(order.get('item-price', 0) or 0)
                    analysis[category]['revenue'] += item_price

                    # Analisi mensile
                    if purchase_date:
                        dt = datetime.fromisoformat(purchase_date.replace('Z', '+00:00'))
                        month_key = dt.strftime('%Y-%m')

                        if is_business:
                            analysis['monthly_business'][month_key]['orders'] += 1
                            analysis['monthly_business'][month_key]['revenue'] += item_price
                        else:
                            analysis['monthly_consumer'][month_key]['orders'] += 1
                            analysis['monthly_consumer'][month_key]['revenue'] += item_price
                except (ValueError, TypeError):
                    pass

        # Ordina dati mensili
        analysis['monthly_business'] = dict(sorted(analysis['monthly_business'].items()))
        analysis['monthly_consumer'] = dict(sorted(analysis['monthly_consumer'].items()))

        return analysis

    def _add_logo(self, slide):
        """Aggiunge il logo Mela alla slide."""
        logo_path = os.path.join(self.data_directory, "mela_logo.png")
        if os.path.exists(logo_path):
            # Posizione in basso a destra, dimensioni più piccole
            left = Inches(9.0)
            top = Inches(5.0)
            width = Inches(0.8)

            try:
                slide.shapes.add_picture(logo_path, left, top, width=width)
            except Exception as e:
                print(f"⚠️  Impossibile aggiungere logo: {e}")

    def create_powerpoint(self, output_filename: str = "Amazon_Orders_Mela_Report.pptx"):
        """Crea una presentazione con stile Mela."""
        if not PPTX_AVAILABLE:
            print("❌ python-pptx non disponibile")
            return

        if len(self.orders) == 0:
            print("❌ Nessun ordine caricato. Impossibile creare la presentazione.")
            return

        print(f"\n📊 Creazione presentazione con stile Mela")
        print(f"   Output: {output_filename}")

        # Estrai logo dal template se non esiste già
        self._extract_logo_from_template()

        # Crea nuova presentazione in formato widescreen (come template Mela)
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

        # Crea le nostre slide personalizzate
        self._create_title_slide(prs)
        self._create_summary_slide(prs)
        self._create_revenue_slide(prs)
        self._create_performance_slide(prs)
        self._create_fba_fbm_slide(prs)
        self._create_business_slide(prs)
        self._create_promotions_slide(prs)
        self._create_geography_slide(prs)
        self._create_market_trends_slide(prs)
        self._create_products_slide(prs)
        self._create_insights_slide(prs)

        # Aggiungi logo Mela a tutte le slide
        for slide in prs.slides:
            self._add_logo(slide)

        # Salva
        prs.save(output_filename)
        print(f"✅ Presentazione salvata: {output_filename}")

    def _extract_logo_from_template(self):
        """Estrae il logo dal template Mela se non esiste già."""
        logo_path = os.path.join(self.data_directory, "mela_logo.png")
        template_path = os.path.join(self.data_directory, "Template Mela - Chiaro.pptx")

        if not os.path.exists(logo_path) and os.path.exists(template_path):
            try:
                template = Presentation(template_path)
                if len(template.slides) > 0:
                    slide = template.slides[0]
                    # Il logo è la prima shape (immagine)
                    for shape in slide.shapes:
                        if shape.shape_type == 13:  # PICTURE
                            image = shape.image
                            with open(logo_path, 'wb') as f:
                                f.write(image.blob)
                            print(f"✓ Logo Mela estratto dal template")
                            break
            except Exception as e:
                print(f"⚠️  Impossibile estrarre logo: {e}")

    def _create_title_slide(self, prs):
        """Crea slide titolo."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

        # Background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(5.625)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.COLORS['bg_light']
        bg.line.fill.background()
        slide.shapes._spTree.remove(bg._element)
        slide.shapes._spTree.insert(2, bg._element)

        # Titolo
        title_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(1))
        tf = title_box.text_frame
        tf.text = "Amazon Orders Analysis"
        p = tf.paragraphs[0]
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['primary']
        p.alignment = PP_ALIGN.CENTER

        # Sottotitolo
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(0.6))
        tf = subtitle_box.text_frame
        tf.text = f"Report Completo • {len(self.orders)} Ordini • {datetime.now().strftime('%B %Y')}"
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.color.rgb = self.COLORS['text_light']
        p.alignment = PP_ALIGN.CENTER

    def _create_summary_slide(self, prs):
        """Crea slide summary."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

        # Titolo
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        tf = title_box.text_frame
        tf.text = "Executive Summary"
        p = tf.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['primary']

        summary = self.get_summary()

        # Metric cards stile Mela
        success_rate = (summary['by_status']['Shipped']/summary['total_orders'])*100 if summary['total_orders'] > 0 else 0
        metrics = [
            ("Totale Ordini", f"{summary['total_orders']}", self.COLORS['primary']),
            ("Revenue", f"€{summary['total_revenue']:,.0f}", self.COLORS['secondary']),
            ("Articoli", f"{summary['total_items_sold']}", self.COLORS['accent1']),
            ("Successo", f"{success_rate:.0f}%", self.COLORS['accent2']),
        ]

        y = 1.8
        for i, (title, value, color) in enumerate(metrics):
            x = 0.8 + (i % 4) * 2.3

            # Card minimalista
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y),
                Inches(2.0), Inches(1.4)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = self.COLORS['white']
            card.line.color.rgb = color
            card.line.width = Pt(2)

            tf = card.text_frame
            tf.clear()
            tf.margin_top = Inches(0.2)
            tf.margin_bottom = Inches(0.1)

            # Titolo
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(13)
            p.font.color.rgb = self.COLORS['text_light']
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(5)

            # Valore
            p2 = tf.add_paragraph()
            p2.text = value
            p2.font.size = Pt(28)
            p2.font.bold = True
            p2.font.color.rgb = color
            p2.alignment = PP_ALIGN.CENTER

        # Info box sotto
        info_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1.5))
        tf = info_box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        shipped = summary['by_status']['Shipped']
        cancelled = summary['by_status']['Cancelled']
        pending = summary['by_status']['Pending']
        p.text = f"Status: {shipped} spediti • {cancelled} cancellati • {pending} in attesa"
        p.font.size = Pt(13)
        p.font.color.rgb = self.COLORS['text_light']
        p.alignment = PP_ALIGN.CENTER

    def _create_revenue_slide(self, prs):
        """Crea slide revenue."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Titolo
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        tf = title_box.text_frame
        tf.text = "Revenue Mensile"
        p = tf.paragraphs[0]
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['primary']

        monthly = self.get_monthly_trends()

        # Grafico
        chart_data = CategoryChartData()
        chart_data.categories = list(monthly.keys())
        chart_data.add_series('Revenue', [monthly[m]['revenue'] for m in monthly.keys()])

        x, y, cx, cy = Inches(0.8), Inches(1.4), Inches(8.4), Inches(3.7)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart

        chart.has_legend = False

        # Titolo grafico
        chart.chart_title.has_text_frame = True
        chart.chart_title.text_frame.text = "Revenue Mensile (Ordini Spediti)"
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = self.COLORS['text_dark']

        # Data labels
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.font.size = Pt(9)
        data_labels.number_format = '€#,##0'

        # Assi - ridimensiona font
        chart.category_axis.tick_labels.font.size = Pt(8)
        chart.value_axis.tick_labels.font.size = Pt(8)

    def _create_performance_slide(self, prs):
        """Crea slide performance."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        tf = title_box.text_frame
        tf.text = "Performance per Mese"
        p = tf.paragraphs[0]
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = self.COLORS["primary"]

        monthly = self.get_monthly_trends()

        chart_data = CategoryChartData()
        chart_data.categories = list(monthly.keys())
        chart_data.add_series('Spediti', [monthly[m]['shipped'] for m in monthly.keys()])
        chart_data.add_series('Cancellati', [monthly[m]['cancelled'] for m in monthly.keys()])

        # Grafico più stretto per fare spazio alla legenda a destra
        x, y, cx, cy = Inches(0.8), Inches(1.4), Inches(7.5), Inches(3.7)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.font.size = Pt(8)
        chart.legend.include_in_layout = False

        # Titolo grafico
        chart.chart_title.has_text_frame = True
        chart.chart_title.text_frame.text = "Confronto Ordini Spediti vs Cancellati"
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = self.COLORS['text_dark']

        # Assi - ridimensiona font
        chart.category_axis.tick_labels.font.size = Pt(8)
        chart.value_axis.tick_labels.font.size = Pt(8)

    def _create_geography_slide(self, prs):
        """Crea slide geografia."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        tf = title_box.text_frame
        tf.text = "Top 10 Paesi"
        p = tf.paragraphs[0]
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = self.COLORS["primary"]

        summary = self.get_summary()
        top_countries = dict(sorted(summary['by_country'].items(), key=lambda x: x[1], reverse=True)[:10])

        chart_data = CategoryChartData()
        chart_data.categories = list(top_countries.keys())
        chart_data.add_series('Ordini', list(top_countries.values()))

        x, y, cx, cy = Inches(0.8), Inches(1.4), Inches(8.4), Inches(3.8)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
        ).chart

        chart.has_legend = False

        # Titolo grafico
        chart.chart_title.has_text_frame = True
        chart.chart_title.text_frame.text = "Ordini per Paese"
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = self.COLORS['text_dark']

        # Data labels
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.font.size = Pt(9)

        # Assi - ridimensiona font
        chart.category_axis.tick_labels.font.size = Pt(8)
        chart.value_axis.tick_labels.font.size = Pt(8)

    def _create_market_trends_slide(self, prs):
        """Crea slide con trend mensili per mercato."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        tf = title_box.text_frame
        tf.text = "Vendite Mensili per Mercato"
        p = tf.paragraphs[0]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.COLORS["primary"]

        market_data = self.get_market_monthly_analysis()

        # Prendi i top 5 mercati per revenue totale
        market_totals = {country: sum(months.values()) for country, months in market_data.items()}
        top_markets = sorted(market_totals.items(), key=lambda x: x[1], reverse=True)[:5]

        if top_markets:
            # Raccogli tutti i mesi
            all_months = set()
            for country, _ in top_markets:
                all_months.update(market_data[country].keys())
            all_months = sorted(all_months)

            chart_data = CategoryChartData()
            chart_data.categories = all_months

            # Aggiungi una serie per ogni mercato top
            colors = [self.COLORS['primary'], self.COLORS['secondary'], self.COLORS['accent1'],
                     self.COLORS['accent2'], self.COLORS['accent3']]

            for i, (country, _) in enumerate(top_markets):
                revenues = [market_data[country].get(month, 0) for month in all_months]
                chart_data.add_series(country, revenues)

            x, y, cx, cy = Inches(0.8), Inches(1.2), Inches(7.5), Inches(4.0)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
            ).chart

            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.RIGHT
            chart.legend.font.size = Pt(8)
            chart.legend.include_in_layout = False

            # Titolo grafico
            chart.chart_title.has_text_frame = True
            chart.chart_title.text_frame.text = "Revenue Mensile - Top 5 Mercati"
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
            chart.chart_title.text_frame.paragraphs[0].font.color.rgb = self.COLORS['text_dark']

            # Assi
            chart.category_axis.tick_labels.font.size = Pt(8)
            chart.value_axis.tick_labels.font.size = Pt(8)

            # Colora le linee
            for i, series in enumerate(chart.series):
                if i < len(colors):
                    series.format.line.color.rgb = colors[i]
                    series.format.line.width = Pt(2.5)

    def _create_fba_fbm_slide(self, prs):
        """Crea slide confronto FBA vs FBM."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        tf = title_box.text_frame
        tf.text = "FBA vs FBM"
        p = tf.paragraphs[0]
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = self.COLORS["primary"]

        analysis = self.get_fba_fbm_analysis()

        # Cards comparative in alto
        y_cards = 1.3

        # Card FBA
        fba_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(y_cards),
            Inches(4.2), Inches(1.3)
        )
        fba_card.fill.solid()
        fba_card.fill.fore_color.rgb = self.COLORS['white']
        fba_card.line.color.rgb = self.COLORS['secondary']
        fba_card.line.width = Pt(3)

        tf = fba_card.text_frame
        tf.clear()
        tf.margin_top = Inches(0.15)
        tf.margin_left = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = "FBA (Fulfillment by Amazon)"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['secondary']
        p.space_after = Pt(8)

        fba_data = analysis['fba']
        cancel_rate_fba = (fba_data['cancelled'] / fba_data['orders'] * 100) if fba_data['orders'] > 0 else 0

        p = tf.add_paragraph()
        p.text = f"Ordini: {fba_data['orders']} • Revenue: €{fba_data['revenue']:,.0f}"
        p.font.size = Pt(12)
        p.font.color.rgb = self.COLORS['text_dark']
        p.space_after = Pt(4)

        p = tf.add_paragraph()
        p.text = f"Articoli: {fba_data['items']} • Cancellazioni: {cancel_rate_fba:.1f}%"
        p.font.size = Pt(11)
        p.font.color.rgb = self.COLORS['text_light']

        # Card FBM
        fbm_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(5.2), Inches(y_cards),
            Inches(4.2), Inches(1.3)
        )
        fbm_card.fill.solid()
        fbm_card.fill.fore_color.rgb = self.COLORS['white']
        fbm_card.line.color.rgb = self.COLORS['accent1']
        fbm_card.line.width = Pt(3)

        tf = fbm_card.text_frame
        tf.clear()
        tf.margin_top = Inches(0.15)
        tf.margin_left = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = "FBM (Fulfillment by Merchant)"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['accent1']
        p.space_after = Pt(8)

        fbm_data = analysis['fbm']
        cancel_rate_fbm = (fbm_data['cancelled'] / fbm_data['orders'] * 100) if fbm_data['orders'] > 0 else 0

        p = tf.add_paragraph()
        p.text = f"Ordini: {fbm_data['orders']} • Revenue: €{fbm_data['revenue']:,.0f}"
        p.font.size = Pt(12)
        p.font.color.rgb = self.COLORS['text_dark']
        p.space_after = Pt(4)

        p = tf.add_paragraph()
        p.text = f"Articoli: {fbm_data['items']} • Cancellazioni: {cancel_rate_fbm:.1f}%"
        p.font.size = Pt(11)
        p.font.color.rgb = self.COLORS['text_light']

        # Grafico confronto mensile revenue
        all_months = sorted(set(list(analysis['monthly_fba'].keys()) + list(analysis['monthly_fbm'].keys())))

        chart_data = CategoryChartData()
        chart_data.categories = all_months

        fba_revenues = [analysis['monthly_fba'].get(m, {'revenue': 0})['revenue'] for m in all_months]
        fbm_revenues = [analysis['monthly_fbm'].get(m, {'revenue': 0})['revenue'] for m in all_months]

        chart_data.add_series('FBA', fba_revenues)
        chart_data.add_series('FBM', fbm_revenues)

        # Grafico più stretto per fare spazio alla legenda a destra
        x, y, cx, cy = Inches(0.8), Inches(2.8), Inches(7.5), Inches(2.5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.font.size = Pt(8)
        chart.legend.include_in_layout = False

        # Titolo grafico
        chart.chart_title.has_text_frame = True
        chart.chart_title.text_frame.text = "Revenue Mensile: FBA vs FBM"
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = self.COLORS['text_dark']

        # Assi - ridimensiona font
        chart.category_axis.tick_labels.font.size = Pt(8)
        chart.value_axis.tick_labels.font.size = Pt(8)

    def _create_products_slide(self, prs):
        """Crea slide prodotti."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        tf = title_box.text_frame
        tf.text = "Top 8 Prodotti (SKU)"
        p = tf.paragraphs[0]
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = self.COLORS["primary"]

        products = self.get_products_summary()[:8]

        chart_data = CategoryChartData()
        # Usa SKU invece dei nomi prodotti
        chart_data.categories = [p['sku'] for p in products]
        chart_data.add_series('Revenue', [p['total_revenue'] for p in products])

        x, y, cx, cy = Inches(0.8), Inches(1.4), Inches(8.4), Inches(3.8)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
        ).chart

        chart.has_legend = False

        # Titolo grafico
        chart.chart_title.has_text_frame = True
        chart.chart_title.text_frame.text = "Revenue per Prodotto"
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = self.COLORS['text_dark']

        # Data labels con formato euro
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.font.size = Pt(9)
        data_labels.number_format = '€#,##0'

        # Assi - ridimensiona font
        chart.category_axis.tick_labels.font.size = Pt(8)
        chart.value_axis.tick_labels.font.size = Pt(8)

    def _create_business_slide(self, prs):
        """Crea slide confronto Business vs Consumer."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        tf = title_box.text_frame
        tf.text = "Ordini Business vs Consumer"
        p = tf.paragraphs[0]
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = self.COLORS["primary"]

        business_data = self.get_business_analysis()

        # Card comparative Business
        y = 1.0
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.0), Inches(y),
            Inches(3.8), Inches(1.8)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = self.COLORS['white']
        card.line.color.rgb = self.COLORS['secondary']
        card.line.width = Pt(3)

        tf = card.text_frame
        tf.clear()
        tf.margin_top = Inches(0.2)
        tf.margin_left = Inches(0.15)
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = "BUSINESS ORDERS"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['secondary']
        p.alignment = PP_ALIGN.CENTER

        p = tf.add_paragraph()
        p.text = f"Ordini: {business_data['business']['orders']}"
        p.font.size = Pt(13)
        p.font.color.rgb = self.COLORS['text_dark']
        p.space_before = Pt(8)

        p = tf.add_paragraph()
        p.text = f"Revenue: €{business_data['business']['revenue']:.2f}"
        p.font.size = Pt(13)
        p.font.color.rgb = self.COLORS['text_dark']

        p = tf.add_paragraph()
        p.text = f"Items: {business_data['business']['items']}"
        p.font.size = Pt(13)
        p.font.color.rgb = self.COLORS['text_dark']

        p = tf.add_paragraph()
        p.text = f"Cancellati: {business_data['business']['cancelled']}"
        p.font.size = Pt(13)
        p.font.color.rgb = self.COLORS['text_dark']

        # Card comparative Consumer
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(5.2), Inches(y),
            Inches(3.8), Inches(1.8)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = self.COLORS['white']
        card.line.color.rgb = self.COLORS['accent1']
        card.line.width = Pt(3)

        tf = card.text_frame
        tf.clear()
        tf.margin_top = Inches(0.2)
        tf.margin_left = Inches(0.15)
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = "CONSUMER ORDERS"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['accent1']
        p.alignment = PP_ALIGN.CENTER

        p = tf.add_paragraph()
        p.text = f"Ordini: {business_data['consumer']['orders']}"
        p.font.size = Pt(13)
        p.font.color.rgb = self.COLORS['text_dark']
        p.space_before = Pt(8)

        p = tf.add_paragraph()
        p.text = f"Revenue: €{business_data['consumer']['revenue']:.2f}"
        p.font.size = Pt(13)
        p.font.color.rgb = self.COLORS['text_dark']

        p = tf.add_paragraph()
        p.text = f"Items: {business_data['consumer']['items']}"
        p.font.size = Pt(13)
        p.font.color.rgb = self.COLORS['text_dark']

        p = tf.add_paragraph()
        p.text = f"Cancellati: {business_data['consumer']['cancelled']}"
        p.font.size = Pt(13)
        p.font.color.rgb = self.COLORS['text_dark']

        # Grafico revenue mensile Business vs Consumer
        if business_data['monthly_business'] or business_data['monthly_consumer']:
            # Unisci tutti i mesi
            all_months = sorted(set(
                list(business_data['monthly_business'].keys()) +
                list(business_data['monthly_consumer'].keys())
            ))

            chart_data = CategoryChartData()
            chart_data.categories = all_months

            business_revenue = [business_data['monthly_business'].get(m, {}).get('revenue', 0) for m in all_months]
            consumer_revenue = [business_data['monthly_consumer'].get(m, {}).get('revenue', 0) for m in all_months]

            chart_data.add_series('Business', business_revenue)
            chart_data.add_series('Consumer', consumer_revenue)

            # Grafico più stretto per fare spazio alla legenda a destra
            x, y, cx, cy = Inches(1.0), Inches(3.1), Inches(7.2), Inches(2.2)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            ).chart

            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.RIGHT
            chart.legend.font.size = Pt(8)
            chart.legend.include_in_layout = False

            # Titolo grafico
            chart.chart_title.has_text_frame = True
            chart.chart_title.text_frame.text = "Revenue Mensile: Business vs Consumer"
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
            chart.chart_title.text_frame.paragraphs[0].font.color.rgb = self.COLORS['text_dark']

            # Colori serie
            chart.series[0].format.fill.solid()
            chart.series[0].format.fill.fore_color.rgb = self.COLORS['secondary']
            chart.series[1].format.fill.solid()
            chart.series[1].format.fill.fore_color.rgb = self.COLORS['accent1']

            # Assi - ridimensiona font
            chart.category_axis.tick_labels.font.size = Pt(8)
            chart.value_axis.tick_labels.font.size = Pt(8)

    def _create_promotions_slide(self, prs):
        """Crea slide analisi promozioni."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        tf = title_box.text_frame
        tf.text = "Analisi Promozioni"
        p = tf.paragraphs[0]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.COLORS["primary"]

        promo_data = self.get_promotions_analysis()

        # Card riassuntive - più compatte e distanziate dai grafici
        total_promo = promo_data['total_with_promotions']
        num_promo_types = len(promo_data['by_promotion_type'])
        num_skus = len(promo_data['skus_on_promotion'])

        cards = [
            ("Ordini Promo", str(total_promo), self.COLORS['accent1']),
            ("Tipi Promo", str(num_promo_types), self.COLORS['accent2']),
            ("SKU Coinvolti", str(num_skus), self.COLORS['accent3']),
        ]

        y = 0.95
        for i, (label, value, color) in enumerate(cards):
            x = 0.8 + (i * 3.0)

            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y),
                Inches(2.7), Inches(0.85)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = self.COLORS['white']
            card.line.color.rgb = color
            card.line.width = Pt(2)

            tf = card.text_frame
            tf.clear()
            tf.margin_top = Inches(0.12)
            tf.margin_left = Inches(0.1)
            tf.word_wrap = True

            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(10)
            p.font.color.rgb = self.COLORS['text_light']
            p.alignment = PP_ALIGN.CENTER

            p = tf.add_paragraph()
            p.text = value
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = color
            p.alignment = PP_ALIGN.CENTER

        # Grafico tipi di promozione - più in basso
        if promo_data['by_promotion_type']:
            promo_types = sorted(
                promo_data['by_promotion_type'].items(),
                key=lambda x: x[1]['count'],
                reverse=True
            )

            chart_data = CategoryChartData()
            # Tronca nomi lunghi delle promozioni
            chart_data.categories = [
                (name[:30] + '...' if len(name) > 30 else name)
                for name, _ in promo_types
            ]
            chart_data.add_series('Ordini', [data['count'] for _, data in promo_types])

            x, y, cx, cy = Inches(0.6), Inches(2.1), Inches(4.4), Inches(3.2)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            ).chart

            chart.has_legend = False

            # Titolo grafico
            chart.chart_title.has_text_frame = True
            chart.chart_title.text_frame.text = "Ordini per Tipo Promozione"
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(13)
            chart.chart_title.text_frame.paragraphs[0].font.color.rgb = self.COLORS['text_dark']

            # Data labels
            chart.plots[0].has_data_labels = True
            data_labels = chart.plots[0].data_labels
            data_labels.font.size = Pt(10)

            # Assi - ridimensiona font
            chart.category_axis.tick_labels.font.size = Pt(8)
            chart.value_axis.tick_labels.font.size = Pt(8)

        # Top 5 SKU in promozione
        if promo_data['skus_on_promotion']:
            top_skus = sorted(
                promo_data['skus_on_promotion'].items(),
                key=lambda x: x[1],
                reverse=True
            )[:5]

            chart_data = CategoryChartData()
            chart_data.categories = [sku for sku, _ in top_skus]
            chart_data.add_series('Ordini', [count for _, count in top_skus])

            x, y, cx, cy = Inches(5.3), Inches(2.1), Inches(4.2), Inches(3.2)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
            ).chart

            chart.has_legend = False

            # Titolo grafico
            chart.chart_title.has_text_frame = True
            chart.chart_title.text_frame.text = "Top 5 SKU in Promozione"
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(13)
            chart.chart_title.text_frame.paragraphs[0].font.color.rgb = self.COLORS['text_dark']

            # Data labels
            chart.plots[0].has_data_labels = True
            data_labels = chart.plots[0].data_labels
            data_labels.font.size = Pt(10)

            # Assi - ridimensiona font
            chart.category_axis.tick_labels.font.size = Pt(8)
            chart.value_axis.tick_labels.font.size = Pt(8)

    def _create_insights_slide(self, prs):
        """Crea slide insights."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        tf = title_box.text_frame
        tf.text = "Key Insights"
        p = tf.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.COLORS["primary"]

        summary = self.get_summary()
        products = self.get_products_summary()
        monthly = self.get_monthly_trends()

        # Calcola metriche avanzate
        best_month = max(monthly.items(), key=lambda x: x[1]['revenue'])
        success_rate = (summary['by_status']['Shipped']/summary['total_orders'])*100

        # Top paese per revenue (non solo per numero ordini)
        country_revenue = defaultdict(float)
        for order in self.orders:
            if order.get('order-status') == 'Shipped':
                country = order.get('ship-country', 'Unknown')
                try:
                    revenue = float(order.get('item-price', 0) or 0)
                    country_revenue[country] += revenue
                except (ValueError, TypeError):
                    pass

        top_country_revenue = max(country_revenue.items(), key=lambda x: x[1]) if country_revenue else ('N/A', 0)

        # Top prodotto (SKU) per revenue
        top_product = products[0] if products else {'sku': 'N/A', 'total_revenue': 0}

        # Valore medio ordine
        avg_order_value = summary['total_revenue'] / summary['by_status']['Shipped'] if summary['by_status']['Shipped'] > 0 else 0

        # Card insights stile Mela - metriche migliorate
        insights = [
            ("✓", "Success Rate", f"{success_rate:.1f}%\nordini spediti", self.COLORS['accent2']),
            ("★", "Best Month", f"{best_month[0]}\n€{best_month[1]['revenue']:.0f}", self.COLORS['primary']),
            ("€", "Top Market", f"{top_country_revenue[0]}\n€{top_country_revenue[1]:.0f} revenue", self.COLORS['secondary']),
            ("◉", "Top SKU", f"{top_product['sku']}\n€{top_product['total_revenue']:.0f}", self.COLORS['accent1']),
            ("Ø", "Avg Order", f"€{avg_order_value:.2f}\nper ordine", self.COLORS['accent3']),
        ]

        # Prima riga: 3 card
        y_row1 = 1.3
        y_row2 = 3.5

        for i, (icon, title, text, color) in enumerate(insights):
            if i < 3:
                # Prima riga
                x = 0.8 + (i % 3) * 3.0
                y = y_row1
            else:
                # Seconda riga - centrate
                x = 2.3 + ((i - 3) % 2) * 3.0
                y = y_row2

            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y),
                Inches(2.7), Inches(1.9)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = self.COLORS['white']
            card.line.color.rgb = color
            card.line.width = Pt(3)

            tf = card.text_frame
            tf.clear()
            tf.margin_top = Inches(0.18)
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            tf.word_wrap = True

            # Icon
            p = tf.paragraphs[0]
            p.text = icon
            p.font.size = Pt(34)
            p.font.color.rgb = color
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(8)

            # Title
            p = tf.add_paragraph()
            p.text = title
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = self.COLORS['text_dark']
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(6)

            # Text - con line spacing per separare righe
            p = tf.add_paragraph()
            p.text = text
            p.font.size = Pt(14)
            p.font.color.rgb = color
            p.alignment = PP_ALIGN.CENTER
            p.line_spacing = 1.2


def main():
    """Funzione principale."""
    print("Amazon Orders Analyzer - Template Mela")
    print("="*70)

    if not PPTX_AVAILABLE:
        print("\n⚠️  Installare python-pptx")
        return

    # Usa la directory dello script, non la directory corrente
    script_dir = os.path.dirname(os.path.abspath(__file__))

    # Inizializza l'analyzer con la directory dello script
    analyzer = AmazonOrdersAnalyzer(script_dir)

    # Carica i file (cerca tutti i file .txt che non siano template o readme)
    analyzer.load_files("*.txt")

    # Crea PowerPoint con template Mela
    analyzer.create_powerpoint()

    print("\n" + "="*70)
    print("✅ Presentazione creata con successo!")
    print("🍎 Usando stile Template Mela")
    print("📊 File: Amazon_Orders_Mela_Report.pptx")
    print("="*70)


if __name__ == "__main__":
    main()
